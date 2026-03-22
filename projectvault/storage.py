"""
ProjectVault Storage Layer

Manages the local filesystem and SQLite database that back all vault operations.
Storage is organized as:

    ~/.projectvault/
        config.json
        projectvault.db        (SQLite with FTS5)
        vaults/
            {vault_id}/
                docs/
                    {doc_id}/
                        current{.ext}
                        metadata.json
                        extracted.txt
                        history/
                            v1{.ext}
                            v2{.ext}
"""

import json
import os
import shutil
import sqlite3
import uuid
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DEFAULT_ROOT = Path.home() / ".projectvault"
CONFIG_FILE = "config.json"
DB_FILE = "projectvault.db"
VAULTS_DIR = "vaults"

# File size limits (bytes)
MAX_FILE_SIZE = 30 * 1024 * 1024  # 30 MB per file (matches Claude Projects)

# Free tier limits
FREE_MAX_VAULTS = 3
FREE_MAX_DOCS_PER_VAULT = 50
FREE_MAX_STORAGE_BYTES = 500 * 1024 * 1024  # 500 MB total
FREE_MAX_VERSIONS = 5


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def extract_text(file_path: Path) -> str:
    """Extract searchable text from a file based on its extension.

    For text-native formats, reads the file directly.
    For rich documents (PDF, DOCX, XLSX, PPTX), uses appropriate libraries.
    For binary/image/audio files, returns empty string (no extraction in v1).
    """
    suffix = file_path.suffix.lower()

    # Text-native formats -- read directly
    if suffix in (".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml",
                   ".py", ".js", ".ts", ".sql", ".yaml", ".yml", ".toml",
                   ".ini", ".cfg", ".conf", ".log", ".sh", ".bat", ".ps1"):
        try:
            return file_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""

    # PDF
    if suffix == ".pdf":
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    # DOCX
    if suffix == ".docx":
        try:
            import docx
            doc = docx.Document(str(file_path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""

    # XLSX
    if suffix == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        text_parts.append("\t".join(cells))
            wb.close()
            return "\n".join(text_parts)
        except Exception:
            return ""

    # PPTX
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                text_parts.append(text)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    # Images, audio, and unknown formats -- no extraction in v1
    return ""


# ---------------------------------------------------------------------------
# Database setup
# ---------------------------------------------------------------------------

def _init_db(db_path: Path) -> None:
    """Create the SQLite database with FTS5 tables if they don't exist."""
    conn = sqlite3.connect(str(db_path))
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")

    conn.executescript("""
        CREATE TABLE IF NOT EXISTS vaults (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            description TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            archived INTEGER DEFAULT 0,
            tags TEXT DEFAULT '[]',
            linked_projects TEXT DEFAULT '[]'
        );

        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            vault_id TEXT NOT NULL,
            name TEXT NOT NULL,
            original_filename TEXT NOT NULL,
            file_extension TEXT DEFAULT '',
            category TEXT DEFAULT 'general',
            priority TEXT DEFAULT 'normal',
            tags TEXT DEFAULT '[]',
            notes TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            file_size_bytes INTEGER DEFAULT 0,
            version_count INTEGER DEFAULT 1,
            deleted INTEGER DEFAULT 0,
            FOREIGN KEY (vault_id) REFERENCES vaults(id) ON DELETE CASCADE
        );

        CREATE TABLE IF NOT EXISTS doc_links (
            id TEXT PRIMARY KEY,
            source_vault_id TEXT NOT NULL,
            source_doc_id TEXT NOT NULL,
            target_vault_id TEXT NOT NULL,
            target_doc_id TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY (source_vault_id) REFERENCES vaults(id),
            FOREIGN KEY (target_vault_id) REFERENCES vaults(id)
        );

        CREATE VIRTUAL TABLE IF NOT EXISTS doc_fts USING fts5(
            doc_id,
            vault_id,
            name,
            content,
            tags,
            notes,
            content_rowid='rowid'
        );

        CREATE INDEX IF NOT EXISTS idx_documents_vault ON documents(vault_id);
        CREATE INDEX IF NOT EXISTS idx_documents_deleted ON documents(deleted);
        CREATE INDEX IF NOT EXISTS idx_documents_category ON documents(category);
    """)

    conn.commit()
    conn.close()


# ---------------------------------------------------------------------------
# Storage class
# ---------------------------------------------------------------------------

class VaultStorage:
    """Manages the local filesystem and SQLite database for ProjectVault."""

    def __init__(self, root: Optional[Path] = None):
        self.root = root or DEFAULT_ROOT
        self.db_path = self.root / DB_FILE
        self.vaults_dir = self.root / VAULTS_DIR

        # Ensure directories exist
        self.root.mkdir(parents=True, exist_ok=True)
        self.vaults_dir.mkdir(parents=True, exist_ok=True)

        # Initialize database
        _init_db(self.db_path)

    @contextmanager
    def _db(self):
        """Context manager for database connections."""
        conn = sqlite3.connect(str(self.db_path))
        conn.row_factory = sqlite3.Row
        conn.execute("PRAGMA foreign_keys=ON")
        try:
            yield conn
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            conn.close()

    def _now(self) -> str:
        """Current UTC timestamp as ISO string."""
        return datetime.now(timezone.utc).isoformat()

    # -------------------------------------------------------------------
    # Vault operations
    # -------------------------------------------------------------------

    def create_vault(self, name: str, description: str = "",
                     tags: Optional[List[str]] = None,
                     linked_projects: Optional[List[str]] = None) -> Dict[str, Any]:
        """Create a new vault. Returns vault metadata dict."""
        vault_id = str(uuid.uuid4())[:12]
        now = self._now()
        tags = tags or []
        linked_projects = linked_projects or []

        vault_dir = self.vaults_dir / vault_id / "docs"
        vault_dir.mkdir(parents=True, exist_ok=True)

        with self._db() as conn:
            conn.execute(
                """INSERT INTO vaults (id, name, description, created_at, updated_at, tags, linked_projects)
                   VALUES (?, ?, ?, ?, ?, ?, ?)""",
                (vault_id, name, description, now, now,
                 json.dumps(tags), json.dumps(linked_projects))
            )

        return {
            "id": vault_id,
            "name": name,
            "description": description,
            "created_at": now,
            "updated_at": now,
            "tags": tags,
            "linked_projects": linked_projects,
            "doc_count": 0,
            "total_size_bytes": 0,
        }

    def list_vaults(self, include_archived: bool = False) -> List[Dict[str, Any]]:
        """List all vaults with summary statistics."""
        with self._db() as conn:
            if include_archived:
                rows = conn.execute("SELECT * FROM vaults ORDER BY updated_at DESC").fetchall()
            else:
                rows = conn.execute(
                    "SELECT * FROM vaults WHERE archived = 0 ORDER BY updated_at DESC"
                ).fetchall()

            result = []
            for row in rows:
                vault_id = row["id"]
                stats = conn.execute(
                    """SELECT COUNT(*) as doc_count, COALESCE(SUM(file_size_bytes), 0) as total_size
                       FROM documents WHERE vault_id = ? AND deleted = 0""",
                    (vault_id,)
                ).fetchone()

                result.append({
                    "id": vault_id,
                    "name": row["name"],
                    "description": row["description"],
                    "created_at": row["created_at"],
                    "updated_at": row["updated_at"],
                    "archived": bool(row["archived"]),
                    "tags": json.loads(row["tags"]),
                    "linked_projects": json.loads(row["linked_projects"]),
                    "doc_count": stats["doc_count"],
                    "total_size_bytes": stats["total_size"],
                })
            return result

    def get_vault(self, vault_id: str) -> Optional[Dict[str, Any]]:
        """Get a single vault by ID. Returns None if not found."""
        with self._db() as conn:
            row = conn.execute("SELECT * FROM vaults WHERE id = ?", (vault_id,)).fetchone()
            if not row:
                return None

            stats = conn.execute(
                """SELECT COUNT(*) as doc_count, COALESCE(SUM(file_size_bytes), 0) as total_size
                   FROM documents WHERE vault_id = ? AND deleted = 0""",
                (vault_id,)
            ).fetchone()

            docs = conn.execute(
                """SELECT id, name, category, tags, updated_at, file_size_bytes
                   FROM documents WHERE vault_id = ? AND deleted = 0
                   ORDER BY updated_at DESC""",
                (vault_id,)
            ).fetchall()

            return {
                "id": row["id"],
                "name": row["name"],
                "description": row["description"],
                "created_at": row["created_at"],
                "updated_at": row["updated_at"],
                "archived": bool(row["archived"]),
                "tags": json.loads(row["tags"]),
                "linked_projects": json.loads(row["linked_projects"]),
                "doc_count": stats["doc_count"],
                "total_size_bytes": stats["total_size"],
                "documents": [
                    {
                        "id": d["id"],
                        "name": d["name"],
                        "category": d["category"],
                        "tags": json.loads(d["tags"]),
                        "updated_at": d["updated_at"],
                        "file_size_bytes": d["file_size_bytes"],
                    }
                    for d in docs
                ],
            }

    def find_vault_by_name(self, name: str) -> Optional[Dict[str, Any]]:
        """Find a vault by name (case-insensitive). Returns first match."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT id FROM vaults WHERE LOWER(name) = LOWER(?) AND archived = 0",
                (name,)
            ).fetchone()
            if row:
                return self.get_vault(row["id"])
            return None

    def archive_vault(self, vault_id: str) -> bool:
        """Archive a vault (soft delete). Returns True if found."""
        with self._db() as conn:
            cursor = conn.execute(
                "UPDATE vaults SET archived = 1, updated_at = ? WHERE id = ?",
                (self._now(), vault_id)
            )
            return cursor.rowcount > 0

    def delete_vault(self, vault_id: str) -> bool:
        """Permanently delete a vault and all its documents. Returns True if found."""
        vault_dir = self.vaults_dir / vault_id
        with self._db() as conn:
            # Remove FTS entries for all docs in this vault
            conn.execute("DELETE FROM doc_fts WHERE vault_id = ?", (vault_id,))
            # Remove doc_links
            conn.execute(
                "DELETE FROM doc_links WHERE source_vault_id = ? OR target_vault_id = ?",
                (vault_id, vault_id)
            )
            # Remove documents
            conn.execute("DELETE FROM documents WHERE vault_id = ?", (vault_id,))
            # Remove vault
            cursor = conn.execute("DELETE FROM vaults WHERE id = ?", (vault_id,))
            found = cursor.rowcount > 0

        if vault_dir.exists():
            shutil.rmtree(vault_dir)

        return found

    def link_project(self, vault_id: str, project_name: str) -> bool:
        """Associate a Claude Project name with a vault."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT linked_projects FROM vaults WHERE id = ?", (vault_id,)
            ).fetchone()
            if not row:
                return False
            projects = json.loads(row["linked_projects"])
            if project_name not in projects:
                projects.append(project_name)
                conn.execute(
                    "UPDATE vaults SET linked_projects = ?, updated_at = ? WHERE id = ?",
                    (json.dumps(projects), self._now(), vault_id)
                )
            return True

    # -------------------------------------------------------------------
    # Document operations
    # -------------------------------------------------------------------

    def add_document(self, vault_id: str, name: str, content: bytes,
                     filename: str, tags: Optional[List[str]] = None,
                     category: str = "general", priority: str = "normal",
                     notes: str = "") -> Optional[Dict[str, Any]]:
        """Add a document to a vault. Content is raw bytes.
        Returns document metadata dict, or None if vault not found.
        """
        # Verify vault exists
        vault = self.get_vault(vault_id)
        if not vault:
            return None

        doc_id = str(uuid.uuid4())[:12]
        now = self._now()
        tags = tags or []
        ext = Path(filename).suffix.lower()

        # Write file to disk
        doc_dir = self.vaults_dir / vault_id / "docs" / doc_id
        doc_dir.mkdir(parents=True, exist_ok=True)
        (doc_dir / "history").mkdir(exist_ok=True)

        current_file = doc_dir / f"current{ext}"
        current_file.write_bytes(content)

        # Extract text for search indexing
        extracted = extract_text(current_file)
        (doc_dir / "extracted.txt").write_text(extracted, encoding="utf-8")

        # Write metadata
        meta = {
            "id": doc_id,
            "vault_id": vault_id,
            "name": name,
            "original_filename": filename,
            "file_extension": ext,
            "category": category,
            "priority": priority,
            "tags": tags,
            "notes": notes,
            "created_at": now,
            "updated_at": now,
            "file_size_bytes": len(content),
            "version_count": 1,
        }
        (doc_dir / "metadata.json").write_text(
            json.dumps(meta, indent=2), encoding="utf-8"
        )

        file_size = len(content)

        # Insert into database
        with self._db() as conn:
            conn.execute(
                """INSERT INTO documents
                   (id, vault_id, name, original_filename, file_extension,
                    category, priority, tags, notes, created_at, updated_at,
                    file_size_bytes, version_count)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (doc_id, vault_id, name, filename, ext, category, priority,
                 json.dumps(tags), notes, now, now, file_size, 1)
            )

            # Index in FTS
            conn.execute(
                """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                   VALUES (?, ?, ?, ?, ?, ?)""",
                (doc_id, vault_id, name, extracted, " ".join(tags), notes)
            )

            # Update vault timestamp
            conn.execute(
                "UPDATE vaults SET updated_at = ? WHERE id = ?", (now, vault_id)
            )

        return meta

    def add_document_from_text(self, vault_id: str, name: str, text_content: str,
                               filename: Optional[str] = None,
                               tags: Optional[List[str]] = None,
                               category: str = "general", priority: str = "normal",
                               notes: str = "") -> Optional[Dict[str, Any]]:
        """Convenience method to add a text/markdown document from a string."""
        if filename is None:
            safe_name = name.replace(" ", "_").lower()
            filename = f"{safe_name}.md"
        content_bytes = text_content.encode("utf-8")
        return self.add_document(
            vault_id, name, content_bytes, filename,
            tags=tags, category=category, priority=priority, notes=notes
        )

    def update_document(self, doc_id: str, content: Optional[bytes] = None,
                        filename: Optional[str] = None,
                        name: Optional[str] = None,
                        tags: Optional[List[str]] = None,
                        category: Optional[str] = None,
                        priority: Optional[str] = None,
                        notes: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """Update a document. Auto-versions the previous file if content changes.
        Returns updated metadata, or None if not found.
        """
        with self._db() as conn:
            row = conn.execute(
                "SELECT * FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None

            vault_id = row["vault_id"]
            now = self._now()
            doc_dir = self.vaults_dir / vault_id / "docs" / doc_id
            ext = row["file_extension"]
            current_file = doc_dir / f"current{ext}"
            version_count = row["version_count"]

            # If content is changing, version the current file first
            if content is not None:
                if current_file.exists():
                    history_file = doc_dir / "history" / f"v{version_count}{ext}"
                    shutil.copy2(current_file, history_file)
                    version_count += 1

                if filename:
                    new_ext = Path(filename).suffix.lower()
                    if new_ext != ext:
                        current_file.unlink(missing_ok=True)
                        ext = new_ext
                        current_file = doc_dir / f"current{ext}"

                current_file.write_bytes(content)

                # Re-extract text
                extracted = extract_text(current_file)
                (doc_dir / "extracted.txt").write_text(extracted, encoding="utf-8")

                file_size = len(content)
            else:
                file_size = row["file_size_bytes"]
                extracted = None

            # Build update query
            updates = {"updated_at": now, "version_count": version_count}
            if name is not None:
                updates["name"] = name
            if tags is not None:
                updates["tags"] = json.dumps(tags)
            if category is not None:
                updates["category"] = category
            if priority is not None:
                updates["priority"] = priority
            if notes is not None:
                updates["notes"] = notes
            if content is not None:
                updates["file_size_bytes"] = file_size
                updates["file_extension"] = ext
                if filename:
                    updates["original_filename"] = filename

            set_clause = ", ".join(f"{k} = ?" for k in updates)
            values = list(updates.values()) + [doc_id]
            conn.execute(
                f"UPDATE documents SET {set_clause} WHERE id = ?", values
            )

            # Update FTS index
            conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (doc_id,))
            final_name = name if name is not None else row["name"]
            final_tags = tags if tags is not None else json.loads(row["tags"])
            final_notes = notes if notes is not None else row["notes"]
            if extracted is None:
                extracted_path = doc_dir / "extracted.txt"
                extracted = extracted_path.read_text(encoding="utf-8") if extracted_path.exists() else ""

            conn.execute(
                """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                   VALUES (?, ?, ?, ?, ?, ?)""",
                (doc_id, vault_id, final_name, extracted,
                 " ".join(final_tags), final_notes)
            )

            # Update vault timestamp
            conn.execute(
                "UPDATE vaults SET updated_at = ? WHERE id = ?", (now, vault_id)
            )

        # Return fresh metadata
        return self.get_document(doc_id)

    def remove_document(self, doc_id: str) -> bool:
        """Soft-delete a document. Returns True if found."""
        with self._db() as conn:
            cursor = conn.execute(
                "UPDATE documents SET deleted = 1, updated_at = ? WHERE id = ? AND deleted = 0",
                (self._now(), doc_id)
            )
            if cursor.rowcount > 0:
                conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (doc_id,))
                return True
            return False

    def get_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """Get document metadata by ID."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT * FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None
            return {
                "id": row["id"],
                "vault_id": row["vault_id"],
                "name": row["name"],
                "original_filename": row["original_filename"],
                "file_extension": row["file_extension"],
                "category": row["category"],
                "priority": row["priority"],
                "tags": json.loads(row["tags"]),
                "notes": row["notes"],
                "created_at": row["created_at"],
                "updated_at": row["updated_at"],
                "file_size_bytes": row["file_size_bytes"],
                "version_count": row["version_count"],
            }

    def get_document_content(self, doc_id: str) -> Optional[str]:
        """Read the extracted text content of a document."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT vault_id FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None
            extracted_path = self.vaults_dir / row["vault_id"] / "docs" / doc_id / "extracted.txt"
            if extracted_path.exists():
                return extracted_path.read_text(encoding="utf-8")
            return ""

    def get_document_raw_path(self, doc_id: str) -> Optional[Path]:
        """Get the path to the raw document file."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT vault_id, file_extension FROM documents WHERE id = ? AND deleted = 0",
                (doc_id,)
            ).fetchone()
            if not row:
                return None
            return self.vaults_dir / row["vault_id"] / "docs" / doc_id / f"current{row['file_extension']}"

    def list_documents(self, vault_id: str, sort_by: str = "updated_at",
                       sort_order: str = "desc", category: Optional[str] = None,
                       tag: Optional[str] = None,
                       limit: int = 50, offset: int = 0) -> Dict[str, Any]:
        """List documents in a vault with filtering and sorting."""
        valid_sorts = {"name", "updated_at", "created_at", "file_size_bytes", "category"}
        if sort_by not in valid_sorts:
            sort_by = "updated_at"
        if sort_order.lower() not in ("asc", "desc"):
            sort_order = "desc"

        with self._db() as conn:
            where_parts = ["vault_id = ?", "deleted = 0"]
            params: list = [vault_id]

            if category:
                where_parts.append("category = ?")
                params.append(category)

            if tag:
                where_parts.append("tags LIKE ?")
                params.append(f'%"{tag}"%')

            where_clause = " AND ".join(where_parts)

            # Get total count
            total = conn.execute(
                f"SELECT COUNT(*) as cnt FROM documents WHERE {where_clause}", params
            ).fetchone()["cnt"]

            # Get page
            rows = conn.execute(
                f"""SELECT * FROM documents WHERE {where_clause}
                    ORDER BY {sort_by} {sort_order}
                    LIMIT ? OFFSET ?""",
                params + [limit, offset]
            ).fetchall()

            docs = []
            for row in rows:
                docs.append({
                    "id": row["id"],
                    "name": row["name"],
                    "original_filename": row["original_filename"],
                    "file_extension": row["file_extension"],
                    "category": row["category"],
                    "priority": row["priority"],
                    "tags": json.loads(row["tags"]),
                    "notes": row["notes"],
                    "updated_at": row["updated_at"],
                    "file_size_bytes": row["file_size_bytes"],
                    "version_count": row["version_count"],
                })

            return {
                "total": total,
                "count": len(docs),
                "offset": offset,
                "has_more": total > offset + len(docs),
                "documents": docs,
            }

    # -------------------------------------------------------------------
    # Search
    # -------------------------------------------------------------------

    def search(self, query: str, vault_id: Optional[str] = None,
               limit: int = 20, offset: int = 0) -> Dict[str, Any]:
        """Full-text search across document contents using FTS5.

        If vault_id is provided, searches only that vault.
        Otherwise searches all vaults (cross-vault search).
        """
        with self._db() as conn:
            if vault_id:
                rows = conn.execute(
                    """SELECT doc_id, vault_id, name,
                              snippet(doc_fts, 3, '>>>', '<<<', '...', 40) as snippet,
                              rank
                       FROM doc_fts
                       WHERE doc_fts MATCH ? AND vault_id = ?
                       ORDER BY rank
                       LIMIT ? OFFSET ?""",
                    (query, vault_id, limit, offset)
                ).fetchall()
            else:
                rows = conn.execute(
                    """SELECT doc_id, vault_id, name,
                              snippet(doc_fts, 3, '>>>', '<<<', '...', 40) as snippet,
                              rank
                       FROM doc_fts
                       WHERE doc_fts MATCH ?
                       ORDER BY rank
                       LIMIT ? OFFSET ?""",
                    (query, limit, offset)
                ).fetchall()

            results = []
            for row in rows:
                # Get vault name for context
                vault_row = conn.execute(
                    "SELECT name FROM vaults WHERE id = ?", (row["vault_id"],)
                ).fetchone()
                vault_name = vault_row["name"] if vault_row else "Unknown"

                results.append({
                    "doc_id": row["doc_id"],
                    "vault_id": row["vault_id"],
                    "vault_name": vault_name,
                    "doc_name": row["name"],
                    "snippet": row["snippet"],
                    "relevance_rank": row["rank"],
                })

            return {
                "query": query,
                "scope": vault_id or "all_vaults",
                "count": len(results),
                "results": results,
            }

    def search_by_tag(self, tag: str, vault_id: Optional[str] = None) -> List[Dict[str, Any]]:
        """Find all documents with a specific tag."""
        with self._db() as conn:
            if vault_id:
                rows = conn.execute(
                    """SELECT d.*, v.name as vault_name FROM documents d
                       JOIN vaults v ON d.vault_id = v.id
                       WHERE d.tags LIKE ? AND d.vault_id = ? AND d.deleted = 0
                       ORDER BY d.updated_at DESC""",
                    (f'%"{tag}"%', vault_id)
                ).fetchall()
            else:
                rows = conn.execute(
                    """SELECT d.*, v.name as vault_name FROM documents d
                       JOIN vaults v ON d.vault_id = v.id
                       WHERE d.tags LIKE ? AND d.deleted = 0
                       ORDER BY d.updated_at DESC""",
                    (f'%"{tag}"%',)
                ).fetchall()

            return [
                {
                    "id": r["id"],
                    "vault_id": r["vault_id"],
                    "vault_name": r["vault_name"],
                    "name": r["name"],
                    "category": r["category"],
                    "tags": json.loads(r["tags"]),
                    "updated_at": r["updated_at"],
                    "file_size_bytes": r["file_size_bytes"],
                }
                for r in rows
            ]

    # -------------------------------------------------------------------
    # Tagging and metadata
    # -------------------------------------------------------------------

    def tag_document(self, doc_id: str, add_tags: Optional[List[str]] = None,
                     remove_tags: Optional[List[str]] = None) -> Optional[List[str]]:
        """Add or remove tags from a document. Returns final tag list, or None if not found."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT tags FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None

            current_tags = set(json.loads(row["tags"]))
            if add_tags:
                current_tags.update(add_tags)
            if remove_tags:
                current_tags -= set(remove_tags)

            final_tags = sorted(current_tags)
            conn.execute(
                "UPDATE documents SET tags = ?, updated_at = ? WHERE id = ?",
                (json.dumps(final_tags), self._now(), doc_id)
            )

            # Update FTS
            fts_row = conn.execute(
                "SELECT rowid, name, content, notes FROM doc_fts WHERE doc_id = ?", (doc_id,)
            ).fetchone()
            if fts_row:
                conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (doc_id,))
                conn.execute(
                    """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                       VALUES (?, (SELECT vault_id FROM documents WHERE id = ?),
                               ?, ?, ?, ?)""",
                    (doc_id, doc_id, fts_row["name"], fts_row["content"],
                     " ".join(final_tags), fts_row["notes"])
                )

            return final_tags

    def bulk_tag(self, doc_ids: List[str], add_tags: Optional[List[str]] = None,
                 remove_tags: Optional[List[str]] = None) -> int:
        """Apply tag changes to multiple documents. Returns count of modified docs."""
        count = 0
        for doc_id in doc_ids:
            result = self.tag_document(doc_id, add_tags=add_tags, remove_tags=remove_tags)
            if result is not None:
                count += 1
        return count

    # -------------------------------------------------------------------
    # Version history
    # -------------------------------------------------------------------

    def get_doc_history(self, doc_id: str) -> Optional[List[Dict[str, Any]]]:
        """Get version history for a document. Returns list of versions."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT vault_id, file_extension, version_count FROM documents WHERE id = ? AND deleted = 0",
                (doc_id,)
            ).fetchone()
            if not row:
                return None

            history_dir = self.vaults_dir / row["vault_id"] / "docs" / doc_id / "history"
            versions = []
            for i in range(1, row["version_count"]):
                version_file = history_dir / f"v{i}{row['file_extension']}"
                if version_file.exists():
                    stat = version_file.stat()
                    versions.append({
                        "version": i,
                        "file_size_bytes": stat.st_size,
                        "modified_at": datetime.fromtimestamp(
                            stat.st_mtime, tz=timezone.utc
                        ).isoformat(),
                    })

            # Current version
            versions.append({
                "version": row["version_count"],
                "file_size_bytes": 0,  # will be filled from db
                "modified_at": self._now(),
                "current": True,
            })

            return versions

    # -------------------------------------------------------------------
    # Cross-vault operations
    # -------------------------------------------------------------------

    def copy_document(self, doc_id: str, target_vault_id: str) -> Optional[Dict[str, Any]]:
        """Copy a document to another vault. Returns new document metadata."""
        source_doc = self.get_document(doc_id)
        if not source_doc:
            return None

        target_vault = self.get_vault(target_vault_id)
        if not target_vault:
            return None

        # Read the raw file
        raw_path = self.get_document_raw_path(doc_id)
        if not raw_path or not raw_path.exists():
            return None

        content = raw_path.read_bytes()
        return self.add_document(
            target_vault_id,
            name=source_doc["name"],
            content=content,
            filename=source_doc["original_filename"],
            tags=source_doc["tags"],
            category=source_doc["category"],
            priority=source_doc["priority"],
            notes=source_doc["notes"],
        )

    def move_document(self, doc_id: str, target_vault_id: str) -> Optional[Dict[str, Any]]:
        """Move a document to another vault. Returns new document metadata."""
        new_doc = self.copy_document(doc_id, target_vault_id)
        if new_doc:
            self.remove_document(doc_id)
        return new_doc

    # -------------------------------------------------------------------
    # Import / Export
    # -------------------------------------------------------------------

    def import_directory(self, vault_id: str, dir_path: Path,
                         tags: Optional[List[str]] = None,
                         category: str = "imported") -> List[Dict[str, Any]]:
        """Bulk import all supported files from a directory into a vault."""
        imported = []
        if not dir_path.is_dir():
            return imported

        for file_path in sorted(dir_path.iterdir()):
            if file_path.is_file() and not file_path.name.startswith("."):
                if file_path.stat().st_size > MAX_FILE_SIZE:
                    continue
                try:
                    content = file_path.read_bytes()
                    doc_name = file_path.stem.replace("_", " ").replace("-", " ").title()
                    result = self.add_document(
                        vault_id, name=doc_name, content=content,
                        filename=file_path.name, tags=tags, category=category
                    )
                    if result:
                        imported.append(result)
                except Exception:
                    continue

        return imported

    def export_vault(self, vault_id: str, output_dir: Path) -> int:
        """Export all current documents from a vault to a directory.
        Returns number of files exported.
        """
        output_dir.mkdir(parents=True, exist_ok=True)
        count = 0

        with self._db() as conn:
            rows = conn.execute(
                "SELECT id, original_filename, file_extension FROM documents WHERE vault_id = ? AND deleted = 0",
                (vault_id,)
            ).fetchall()

            for row in rows:
                raw_path = self.vaults_dir / vault_id / "docs" / row["id"] / f"current{row['file_extension']}"
                if raw_path.exists():
                    dest = output_dir / row["original_filename"]
                    # Avoid filename collisions
                    counter = 1
                    while dest.exists():
                        stem = Path(row["original_filename"]).stem
                        dest = output_dir / f"{stem}_{counter}{row['file_extension']}"
                        counter += 1
                    shutil.copy2(raw_path, dest)
                    count += 1

        return count
